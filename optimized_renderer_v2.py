# -*- coding: utf-8 -*-
"""
优化模板渲染器 V2 - 每人独立渲染后合并
采用新策略：每个人使用独立的模板副本，填充后合并所有幻灯片
"""
import os
import math
from datetime import datetime
from typing import List, Dict
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt

# 导入DataFormatter
from ppt_renderer import DataFormatter, PaperItem


class OptimizedTemplateRendererV2:
    """优化模板渲染器 V2 - 每人独立渲染后合并"""
    
    def __init__(self, template_path: str):
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"模板文件不存在: {template_path}")
        
        self.template_path = template_path
        self.formatter = DataFormatter()
        print("✅ 渲染器V2创建成功\n")
    
    def render_all(self, resumes: List[Dict], output_dir: str = None) -> str:
        """
        渲染多份简历到一个PPT
        策略：每人使用独立模板，最后合并
        
        Args:
            resumes: 简历数据列表
            output_dir: 输出目录
            
        Returns:
            输出文件路径
        """
        if not resumes:
            raise ValueError("简历列表不能为空")
        
        if output_dir is None:
            output_dir = os.path.dirname(__file__)
        
        print(f"🎨 开始渲染 {len(resumes)} 份简历\n")
        print("📋 策略：每人使用独立模板，最后合并所有幻灯片\n")
        
        # 创建最终PPT，使用第一个模板
        final_prs = Presentation(self.template_path)
        
        # 更新首页人数
        self._update_cover_page(final_prs, len(resumes))
        
        # 为第一份简历填充模板的现有幻灯片（第3-6页）
        print(f"📝 第 1/{len(resumes)} 份简历: {resumes[0].get('姓名', '未知')}")
        self._render_single_resume(final_prs, resumes[0], use_existing_slides=True)
        print()
        
        # 为后续简历创建独立模板并合并
        for i, resume in enumerate(resumes[1:], start=2):
            name = resume.get("姓名", "未知")
            print(f"📝 第 {i}/{len(resumes)} 份简历: {name}")
            
            # 加载独立的模板副本
            temp_prs = Presentation(self.template_path)
            
            # 在独立模板中渲染这份简历
            self._render_single_resume(temp_prs, resume, use_existing_slides=True)
            
            # 将第3-6页的幻灯片复制到最终PPT
            slides_to_copy = []
            for slide_idx in range(2, min(7, len(temp_prs.slides))):  # 第3-7页（可能有额外论文页）
                slides_to_copy.append(temp_prs.slides[slide_idx])
            
            # 检查是否有额外的论文页（第7页之后）
            if len(temp_prs.slides) > 7:
                for slide_idx in range(7, len(temp_prs.slides)):
                    slides_to_copy.append(temp_prs.slides[slide_idx])
            
            print(f"   合并 {len(slides_to_copy)} 页到最终PPT")
            for slide in slides_to_copy:
                self._copy_slide_to_presentation(final_prs, slide)
            
            print()
        
        # 生成输出文件
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(output_dir, f"人才引进简历汇总_{len(resumes)}人_{timestamp}.pptx")
        
        final_prs.save(output_path)
        print(f"✅ PPT已保存: {output_path}")
        print(f"📊 总计幻灯片: {len(final_prs.slides)} 页\n")
        
        return output_path
    
    def _update_cover_page(self, prs: Presentation, total_resumes: int):
        """更新首页人数"""
        if len(prs.slides) > 0:
            for shape in prs.slides[0].shapes:
                if hasattr(shape, 'text_frame'):
                    text = shape.text_frame.text.strip()
                    if "人才引进" in text and "人" in text:
                        import re
                        new_text = re.sub(r'（\d+人）', f'（{total_resumes}人）', text)
                        shape.text_frame.text = new_text
                        print(f"✅ 更新首页人数: {total_resumes}人\n")
                        break
    
    def _render_single_resume(self, prs: Presentation, resume: Dict, use_existing_slides: bool = True):
        """
        在给定的PPT中渲染单份简历
        策略：先判断需要几页，复制好模板页，再统一填充
        
        Args:
            prs: 演示文稿对象
            resume: 简历数据
            use_existing_slides: 是否使用现有幻灯片（True=修改第3-6页，False=添加新幻灯片）
        """
        name = resume.get("姓名", "未知")
        papers_data = resume.get("发表论文情况", [])
        paper_items = self.formatter.extract_paper_items(papers_data)
        
        # 1. 先计算需要的论文页数
        paper_pages = self._calculate_paper_pages(len(paper_items))
        print(f"   📄 论文 {len(paper_items)} 篇，需要 {paper_pages} 页")
        
        # 2. 计算需要的项目页数
        projects = resume.get("获批项目情况", [])
        awards = resume.get("获奖情况", [])
        other = resume.get("其他成果") or []
        if not isinstance(other, list):
            other = []
        
        project_pages = self._calculate_project_pages(projects, awards, other)
        print(f"   📋 项目 {len(projects)} 个，成果 {len(other)} 项，需要 {project_pages} 页")
        
        # 3. 填充第3页：基本情况
        if len(prs.slides) > 2:
            print(f"   填充基本情况页（第3页）")
            self._fill_basic_info_slide(prs.slides[2], resume)
        
        # 4. 处理论文页
        # ⚠️ 重要：模板有2页论文（第4、5页），如果只需要1页或0页，需要删除多余的
        if paper_pages == 0:
            # 删除第4、5页（论文页）
            print(f"   没有论文，删除第4、5页论文模板页")
            if len(prs.slides) > 4:
                rId = prs.slides._sldIdLst[4].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[4]
            if len(prs.slides) > 3:
                rId = prs.slides._sldIdLst[3].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[3]
        elif paper_pages == 1:
            # 只需要1页，删除第5页（索引4）
            print(f"   论文只需1页，删除第5页多余的论文模板页")
            if len(prs.slides) > 4:
                rId = prs.slides._sldIdLst[4].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[4]
        elif paper_pages > 2:
            # 需要超过2页，插入额外页
            extra_paper_pages = paper_pages - 2
            print(f"   需要额外 {extra_paper_pages} 页论文，在第5页后插入模板页")
            
            template_prs = Presentation(self.template_path)
            
            # 在第6页（项目页，索引5）之前插入额外论文页
            insert_position = 5  # 在第6页之前插入
            for i in range(extra_paper_pages):
                print(f"      在位置{insert_position}插入第 {i+1}/{extra_paper_pages} 页额外论文页")
                source_slide = template_prs.slides[4]  # 复制第5页（11行表格）
                self._insert_slide_at_position(prs, source_slide, insert_position)
                insert_position += 1  # 下一次插入的位置向后移
        
        # 5. 填充所有论文页（第4,5,6,7...页）
        if paper_pages > 0:
            paper_offset = 0
            capacities = [8, 11, 11, 11, 11, 11]  # 第1页8行，后续每页11行
            
            for page_num in range(paper_pages):
                slide_idx = 3 + page_num  # 第4页开始（索引3）
                
                if slide_idx < len(prs.slides):
                    capacity = capacities[page_num] if page_num < len(capacities) else 11
                    end = min(paper_offset + capacity, len(paper_items))
                    page_papers = paper_items[paper_offset:end]
                    
                    print(f"   填充论文页 {page_num + 1}/{paper_pages} (第{slide_idx+1}页，{len(page_papers)}/{capacity}篇)")
                    self._fill_paper_slide(prs.slides[slide_idx], resume, page_papers, 
                                         page_num + 1, paper_pages)
                    paper_offset = end
        
        # 6. 处理项目和成果页
        # 新模板结构：第6页=获批项目，第7页=其他成果（已分离）
        # 模板已包含1页项目+1页成果，需要根据实际数量增删
        
        # 计算项目页数和成果页数
        project_count = len(projects)
        project_pages_count = math.ceil(project_count / 10) if project_count > 0 else 1  # 至少1页，每页10个
        other_pages_count = project_pages - project_pages_count
        
        print(f"   项目页数: {project_pages_count}, 成果页数: {other_pages_count}")
        
        # 6.1 处理项目页（模板已有1页）
        if project_pages_count > 1:
            # 需要额外的项目页
            extra_project_pages = project_pages_count - 1
            print(f"   需要额外 {extra_project_pages} 页项目，在第{3 + paper_pages + 1}页后插入")
            
            template_prs = Presentation(self.template_path)
            
            # 项目页的插入位置：基本情况(1) + 论文(N) + 第一个项目页(1)
            project_insert_position = 3 + paper_pages + 1
            for i in range(extra_project_pages):
                print(f"      在位置{project_insert_position}插入第 {i+1}/{extra_project_pages} 页额外项目页")
                source_slide = template_prs.slides[5]  # 复制第6页（项目页）
                self._insert_slide_at_position(prs, source_slide, project_insert_position)
                project_insert_position += 1
        
        # 6.2 处理成果页（模板已有1页）
        if other_pages_count > 1:
            # 需要额外的成果页
            extra_other_pages = other_pages_count - 1
            print(f"   需要额外 {extra_other_pages} 页成果，在第{3 + paper_pages + project_pages_count + 1}页后插入")
            
            template_prs = Presentation(self.template_path)
            
            # 成果页的插入位置：基本情况(1) + 论文(N) + 项目(M) + 第一个成果页(1)
            other_insert_position = 3 + paper_pages + project_pages_count + 1
            for i in range(extra_other_pages):
                print(f"      在位置{other_insert_position}插入第 {i+1}/{extra_other_pages} 页额外成果页")
                source_slide = template_prs.slides[6]  # 复制第7页（成果页）
                self._insert_slide_at_position(prs, source_slide, other_insert_position)
                other_insert_position += 1
        elif other_pages_count == 0:
            # 不需要成果页，删除模板的成果页（第7页）
            # 索引位置：基本(1) + 封面(1) + 模板封面(1) + 基本情况(1) + 论文(N) + 项目(M) = 3 + N + M
            other_template_idx = 3 + paper_pages + project_pages_count
            print(f"   不需要成果页，删除模板成果页（第{other_template_idx+1}页）")
            if other_template_idx < len(prs.slides):
                self._delete_slide(prs, other_template_idx)
        
        # 7. 填充所有项目和成果页
        first_project_idx = 3 + paper_pages  # 第4页 + paper_pages
        for page_num in range(project_pages):
            project_slide_idx = first_project_idx + page_num
            if project_slide_idx < len(prs.slides):
                print(f"   填充项目/成果页 {page_num + 1}/{project_pages} (第{project_slide_idx+1}页)")
                self._fill_project_slide_paginated(
                    prs.slides[project_slide_idx], 
                    resume, 
                    page_num, 
                    project_pages,
                    project_pages_count
                )
        
        # 8. 填充人才办意见页的name（如果存在）
        talent_office_idx = first_project_idx + project_pages  # 所有项目+成果页后面
        if talent_office_idx < len(prs.slides):
            print(f"   填充人才办意见页的姓名（第{talent_office_idx+1}页）")
            self._fill_talent_office_slide(prs.slides[talent_office_idx], resume)
    
    def _calculate_paper_pages(self, total_papers: int) -> int:
        """计算需要的论文页数"""
        if total_papers == 0:
            return 0
        if total_papers <= 8:
            return 1
        remaining = total_papers - 8
        return 1 + math.ceil(remaining / 11)
    
    def _fill_basic_info_slide(self, slide, resume: Dict):
        """填充基本情况页（第3页）"""
        name = resume.get("姓名", "未知")
        
        # 1. 替换所有文本框中的姓名
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                
                # 如果包含"拟聘岗位"，清除后面的具体职位，但保留"拟聘岗位："
                if "拟聘岗位" in text:
                    # 替换为：姓名 + "  拟聘岗位："
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = f"{name}  拟聘岗位："
                    if p.runs:
                        p.runs[0].font.size = Pt(24)
                    continue
                
                # 正常替换姓名
                self._replace_text_preserve_format(shape, "name", name)
                self._replace_text_preserve_format(shape, "苑津山", name)
                
                # 清除模板默认文本
                if "1999" in text and "博士毕业生" in text:
                    continue
                elif "1999" in text:
                    shape.text_frame.clear()
        
        # 2. 查找并填充基本信息文本框
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                if "博士毕业生" in text or "未婚" in text or "周岁" in text or "1999" in text:
                    self._fill_basic_info_textbox(shape, resume)
                    break
        
        # 3. 填充表格（学院编制情况）
        tables = self._find_tables(slide)
        if len(tables) >= 1:
            self._clear_table_data(tables[0], keep_header=True)
    
    def _fill_basic_info_textbox(self, shape, resume: Dict):
        """填充基本信息文本框（优化排版和字体）"""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        
        # 构建基本信息文本
        gender = resume.get("性别", "未知")
        birth_year = resume.get("出生日期", "未知")
        birth_month = "03"  # 默认值
        degree = resume.get("最高学历", "博士")
        marriage = resume.get("婚配情况") or "未婚"
        age = resume.get("年龄", "未知")
        
        # 清空文本框
        shape.text_frame.clear()
        
        # 第一行：基本信息（13pt，醒目）
        p1 = shape.text_frame.paragraphs[0]
        info_line = f"{gender}，{birth_year}年{birth_month}月生，{degree}毕业生，{marriage}，{age}周岁"
        p1.text = info_line
        p1.alignment = PP_ALIGN.LEFT
        for run in p1.runs:
            run.font.size = Pt(13)  # 13pt适中
        
        # 添加空行（改善可读性）
        p_space = shape.text_frame.add_paragraph()
        p_space.text = ""
        
        # 教育经历
        schools = resume.get("就读院校", [])
        if schools:
            # 表头（11pt，加粗）
            p_header = shape.text_frame.add_paragraph()
            p_header.text = "时间\t\t\t院校\t\t\t\t专业\t\t\t学位"
            p_header.alignment = PP_ALIGN.LEFT
            for run in p_header.runs:
                run.font.size = Pt(11)
                run.font.bold = True
            
            # 每条教育经历（10pt）
            for school in schools:
                p = shape.text_frame.add_paragraph()
                time_range = school.get("时间区间", "")
                school_name = school.get("院校", "")
                major = school.get("专业", "")
                degree_val = school.get("学位", "") or ""
                
                # 格式化，确保对齐
                p.text = f"{time_range}\t\t{school_name}\t\t{major}\t\t{degree_val}"
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(10)  # 稍微小一点，更紧凑
    
    def _fill_paper_slide(self, slide, resume: Dict, paper_items: List[PaperItem], 
                         current_page: int, total_pages: int):
        """填充论文页"""
        name = resume.get("姓名", "未知")
        
        # 替换页眉的姓名，保留"拟聘岗位："但不填职位
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                
                # 如果包含"拟聘岗位"，保留但清除具体职位
                if "拟聘岗位" in text:
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = f"{name}  拟聘岗位："
                    if p.runs:
                        p.runs[0].font.size = Pt(24)
                    continue
                
                # 正常替换姓名
                self._replace_text_preserve_format(shape, "name", name)
                self._replace_text_preserve_format(shape, "苑津山", name)
        
        # 添加论文统计信息
        self._add_paper_statistics(slide, resume, current_page, total_pages)
        
        # 填充论文表格
        tables = self._find_tables(slide)
        if tables:
            table = tables[0]
            self._fill_paper_table(table, paper_items)
    
    def _add_paper_statistics(self, slide, resume: Dict, current_page: int, total_pages: int):
        """在论文页添加统计信息"""
        from pptx.util import Pt
        
        # 统计论文数量
        papers_data = resume.get("发表论文情况", [])
        
        # 统计各类论文
        total_papers = 0
        sci_1 = 0
        ccf_a = 0
        
        for paper_group in papers_data:
            titles = paper_group.get("论文题目列表", [])
            count = len(titles)
            total_papers += count
            
            category = paper_group.get("类别", "") or ""
            if "SCI 1区" in category or "SCI1区" in category:
                sci_1 += count
            elif "CCF A" in category:
                ccf_a += count
        
        # 构建统计文本
        stats_text = f"发表论文情况：共发表论文 {total_papers} 篇 | SCI 1区 {sci_1} 篇 | CCF A类 {ccf_a} 篇 (第{current_page}/{total_pages}页，共{total_papers}篇)"
        
        # 查找并更新统计文本框（通常在表格上方）
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text or ""
                # 找到包含"发表论文情况"的文本框
                if "发表论文情况" in text and "共发表论文" not in text:
                    # 清空并设置新文本
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = stats_text
                    # 设置字体
                    for run in p.runs:
                        run.font.size = Pt(11)
                        run.font.bold = False
                    break
    
    def _calculate_project_pages(self, projects: List[Dict], awards: List[Dict], other: List[Dict]) -> int:
        """
        计算需要的项目+成果页数（模板已分离，新版本）
        
        新模板结构：
        - 第6页：获批项目（10行容量）
        - 第7页：其他成果（7行容量，包含获奖情况，每个奖项一行）
        
        分页逻辑：
        - 项目：每页10个，动态分页
        - 成果：获奖每个奖项占1行，其他成果各占1行
          第1页：最多7行（获奖+其他成果）
          后续页：每页7行
        
        Args:
            projects: 项目列表
            awards: 获奖列表  
            other: 其他成果列表
            
        Returns:
            需要的总页数（项目页数 + 成果页数）
        """
        # 计算项目页数（每页10个）
        project_pages = math.ceil(len(projects) / 10) if projects else 1  # 至少1页
        
        # 计算成果页数（获奖每项占1行）
        other_pages = 1  # 默认至少1页
        if awards or other:
            # 总行数 = 获奖数量 + 其他成果数量
            total_rows = len(awards) + len(other)
            
            if total_rows == 0:
                other_pages = 1
            elif total_rows <= 7:
                # 第1页可以放下
                other_pages = 1
            else:
                # 第1页：7行
                # 后续页：每页7行
                remaining = total_rows - 7
                other_pages = 1 + math.ceil(remaining / 7)
        
        total_pages = project_pages + other_pages
        
        return total_pages
    
    def _fill_project_slide_paginated(self, slide, resume: Dict, page_num: int, total_pages: int, project_pages_count: int):
        """
        填充项目/成果页（新模板：项目和成果已分离）
        
        新模板结构：
        - 第6页：获批项目（5行容量）
        - 第7页：其他成果（5行容量）
        
        逻辑：
        - 前N页：项目页（每页5个）
        - 后M页：成果页（第1页：获奖+4个成果，后续：每页5个成果）
        
        Args:
            slide: 幻灯片对象
            resume: 简历数据
            page_num: 当前页码（从0开始）
            total_pages: 总页数
            project_pages_count: 项目页数
        """
        name = resume.get("姓名", "未知")
        
        # 替换页眉的姓名，保留"拟聘岗位："但不填职位
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                
                # 如果包含"拟聘岗位"，保留但清除具体职位
                if "拟聘岗位" in text:
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = f"{name}  拟聘岗位："
                    if p.runs:
                        p.runs[0].font.size = Pt(24)
                    continue
                
                # 正常替换姓名
                self._replace_text_preserve_format(shape, "name", name)
                self._replace_text_preserve_format(shape, "苑津山", name)
        
        projects = resume.get("获批项目情况", [])
        awards = resume.get("获奖情况", [])
        other = resume.get("其他成果") or []
        if not isinstance(other, list):
            other = []
        
        # 判断当前是项目页还是成果页
        is_project_page = page_num < project_pages_count
        
        # 查找表格
        tables = self._find_tables(slide)
        
        if is_project_page:
            # 项目页：填充项目表格（每页10个）
            project_capacity = 10
            start_idx = page_num * project_capacity
            end_idx = min(start_idx + project_capacity, len(projects))
            page_projects = projects[start_idx:end_idx]
            
            if len(tables) >= 1:
                self._fill_projects_table(tables[0], page_projects)
                print(f"      填充 {len(page_projects)} 个项目")
        else:
            # 成果页：填充成果表格（获奖每项一行）
            other_page_num = page_num - project_pages_count  # 成果页的页码（从0开始）
            
            if len(tables) >= 1:
                # 计算当前页应该显示的数据
                awards = resume.get("获奖情况", [])
                total_awards = len(awards)
                total_other = len(other)
                total_items = total_awards + total_other  # 总行数
                
                # 每页7行
                start_idx = other_page_num * 7
                end_idx = min(start_idx + 7, total_items)
                
                # 判断当前页是否包含获奖
                if start_idx < total_awards:
                    # 当前页包含获奖
                    awards_on_page = awards[start_idx:min(end_idx, total_awards)]
                    
                    # 如果还有剩余空间，填充其他成果
                    remaining_space = 7 - len(awards_on_page)
                    if remaining_space > 0 and end_idx > total_awards:
                        other_start = max(0, start_idx - total_awards)
                        other_end = end_idx - total_awards
                        other_on_page = other[other_start:other_end]
                    else:
                        other_on_page = []
                    
                    # 构造resume数据（包含当前页的获奖）
                    page_resume = {"获奖情况": awards_on_page}
                    self._fill_other_achievements_table(
                        tables[0], 
                        other_on_page, 
                        page_resume,
                        show_note=(other_page_num == 0)  # 第1页显示备注
                    )
                    print(f"      填充 {len(awards_on_page)} 个获奖 + {len(other_on_page)} 项成果")
                else:
                    # 当前页只有其他成果
                    other_start = start_idx - total_awards
                    other_end = end_idx - total_awards
                    other_on_page = other[other_start:other_end]
                    
                    self._fill_other_achievements_table(
                        tables[0], 
                        other_on_page, 
                        {"获奖情况": []},  # 不显示获奖
                        show_note=False
                    )
                    print(f"      填充 {len(other_on_page)} 项成果")
                
                # 删除备注文本框（除了第1页）
                if other_page_num > 0:
                    self._remove_note_textbox(slide)
    
    def _fill_talent_office_slide(self, slide, resume: Dict):
        """填充人才办意见页的姓名"""
        name = resume.get("姓名", "未知")
        
        # 替换所有的name占位符（包括大小写变体）
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                # 尝试多种替换方式
                text = shape.text_frame.text
                if "name" in text.lower():
                    # 直接替换整个文本框内容
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run_text = run.text
                            # 替换各种可能的name变体
                            if "name" in run_text.lower():
                                run.text = run_text.replace("name", name).replace("Name", name).replace("NAME", name)
                            if "苑津山" in run_text:
                                run.text = run_text.replace("苑津山", name)
                
                # 也使用原有的替换方法
                self._replace_text_preserve_format(shape, "name", name)
                self._replace_text_preserve_format(shape, "Name", name)
                self._replace_text_preserve_format(shape, "NAME", name)
                self._replace_text_preserve_format(shape, "苑津山", name)
    
    def _fill_basic_info_table(self, table, resume: Dict):
        """填充基本信息表格"""
        info_mapping = {
            0: ("性别", resume.get("性别", "")),
            1: ("出生日期", str(resume.get("出生日期", ""))),
            2: ("婚配情况", resume.get("婚配情况", "") or ""),
            3: ("最高学历", resume.get("最高学历", "")),
        }
        
        for row_idx, (label, value) in info_mapping.items():
            if row_idx < len(table.rows):
                row = table.rows[row_idx]
                if len(row.cells) >= 2:
                    row.cells[1].text = str(value)
    
    def _fill_education_table(self, table, resume: Dict):
        """填充教育经历表格"""
        schools = resume.get("就读院校", [])
        
        for i, school in enumerate(schools):
            if i + 1 >= len(table.rows):
                break
            
            row = table.rows[i + 1]
            if len(row.cells) >= 4:
                row.cells[0].text = school.get("时间区间", "")
                row.cells[1].text = school.get("院校", "")
                row.cells[2].text = school.get("专业", "")
                row.cells[3].text = school.get("学位", "") or ""
    
    def _fill_paper_table(self, table, paper_items: List[PaperItem]):
        """填充论文表格（先清空再填充，并设置文本换行）"""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        
        # 1. 先清空所有数据行（保留表头）
        for i in range(1, len(table.rows)):
            for cell in table.rows[i].cells:
                cell.text = ""
        
        # 2. 填充新数据
        for i, paper in enumerate(paper_items):
            if i + 1 >= len(table.rows):
                break
            
            row = table.rows[i + 1]
            if len(row.cells) >= 5:
                # 设置单元格内容
                row.cells[0].text = paper.journal[:60] if len(paper.journal) > 60 else paper.journal  # 限制期刊名长度
                row.cells[1].text = paper.title[:80] if len(paper.title) > 80 else paper.title  # 限制标题长度
                row.cells[2].text = "1"
                row.cells[3].text = ""
                row.cells[4].text = paper.category[:20] if paper.category and len(paper.category) > 20 else (paper.category or "")
                
                # 设置每个单元格的属性
                for j, cell in enumerate(row.cells):
                    # 启用自动换行
                    cell.text_frame.word_wrap = True
                    
                    # 设置文本属性
                    if cell.text_frame.paragraphs:
                        for paragraph in cell.text_frame.paragraphs:
                            # 垂直对齐
                            paragraph.alignment = PP_ALIGN.LEFT
                            
                            # 设置字体大小
                            for run in paragraph.runs:
                                if run.font.size is None or run.font.size > Pt(9):
                                    run.font.size = Pt(9)  # 稍微减小字体，避免溢出
    
    def _fill_projects_table(self, table, projects: List[Dict]):
        """填充项目表格（先清空再填充，设置自动换行，避免溢出）"""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        
        # 1. 先清空所有数据行（保留表头）
        for i in range(1, len(table.rows)):
            for cell in table.rows[i].cells:
                cell.text = ""
        
        # 2. 填充新数据
        for i, project in enumerate(projects):
            if i + 1 >= len(table.rows):
                break
            
            row = table.rows[i + 1]
            if len(row.cells) >= 5:
                # 项目类别
                row.cells[0].text = project.get("项目类别", "")[:30]  # 限制长度
                
                # 项目名称（可能很长）
                names = project.get("项目名称列表", [])
                names_text = "、".join(names) if names else ""
                row.cells[1].text = names_text[:100] if len(names_text) > 100 else names_text  # 限制长度
                
                row.cells[2].text = str(project.get("项数", ""))
                row.cells[3].text = project.get("年份", "")[:30]  # 限制年份长度
                row.cells[4].text = project.get("备注", "")[:50]  # 限制备注长度
                
                # 设置自动换行和字体
                for cell in row.cells:
                    cell.text_frame.word_wrap = True
                    if cell.text_frame.paragraphs:
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.LEFT
                            for run in paragraph.runs:
                                if run.font.size is None or run.font.size > Pt(10):
                                    run.font.size = Pt(10)
    
    def _fill_other_achievements_table(self, table, other_list: List[Dict], resume: Dict, show_note: bool = True):
        """填充其他成果表格（先清空再填充，设置自动换行）"""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        
        # 1. 先清空所有数据行（保留表头）
        for i in range(1, len(table.rows)):
            for cell in table.rows[i].cells:
                cell.text = ""
        
        # 2. 添加获奖情况（每个奖项一行）
        awards = resume.get("获奖情况", [])
        current_row = 1  # 从第2行开始（索引1）
        
        if awards:
            for award in awards:
                if current_row >= len(table.rows):
                    break  # 表格行数不够，停止
                
                row = table.rows[current_row]
                if len(row.cells) >= 4:
                    # 第一个奖项显示"获奖情况"，后续奖项留空
                    row.cells[0].text = "获奖情况" if current_row == 1 else ""
                    
                    # 奖项名称
                    award_name = award.get("奖项名称", "")
                    row.cells[1].text = award_name[:80] if len(award_name) > 80 else award_name
                    
                    # 数量统计（对于获奖，显示为1项，或者在最后一行显示总数）
                    if current_row == 1:
                        row.cells[2].text = f"{len(awards)}项"
                    else:
                        row.cells[2].text = "1项"
                    
                    # 年份
                    year = award.get("年份", "")
                    row.cells[3].text = str(year) if year else ""
                    
                    # 设置自动换行
                    for cell in row.cells:
                        cell.text_frame.word_wrap = True
                        if cell.text_frame.paragraphs:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.font.size is None or run.font.size > Pt(10):
                                        run.font.size = Pt(10)
                
                current_row += 1
        
        # 3. 添加其他成果
        for i, item in enumerate(other_list):
            row_idx = current_row + i  # 从获奖后面开始
            if row_idx >= len(table.rows):
                break
            
            row = table.rows[row_idx]
            if len(row.cells) >= 4:
                row.cells[0].text = item.get("类别", "")[:20]
                
                names = item.get("名称列表", [])
                names_text = "、".join(names) if names else ""
                row.cells[1].text = names_text[:80] if len(names_text) > 80 else names_text
                
                row.cells[2].text = str(item.get("项数", ""))
                
                year_note = f"{item.get('年份', '')} {item.get('备注', '')}"
                row.cells[3].text = year_note[:50] if len(year_note) > 50 else year_note
                
                # 设置自动换行
                for cell in row.cells:
                    cell.text_frame.word_wrap = True
                    if cell.text_frame.paragraphs:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size is None or run.font.size > Pt(10):
                                    run.font.size = Pt(10)
    
    def _remove_note_textbox(self, slide):
        """删除备注文本框"""
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                if "备注：" in text or "备注:" in text:
                    # 记录需要删除的shape
                    shapes_to_remove.append(shape)
        
        # 删除shape
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
    
    def _clear_table_data(self, table, keep_header: bool = True):
        """清空表格数据"""
        start_row = 1 if keep_header else 0
        for i in range(start_row, len(table.rows)):
            for cell in table.rows[i].cells:
                cell.text = ""
    
    def _replace_text_preserve_format(self, shape, old_text: str, new_text: str):
        """替换文本同时保留格式"""
        if not hasattr(shape, 'text_frame'):
            return
        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
    
    def _find_tables(self, slide):
        """查找幻灯片中的所有表格"""
        tables = []
        for shape in slide.shapes:
            if shape.shape_type == 19:  # GraphicFrame (表格)
                if hasattr(shape, 'table'):
                    tables.append(shape.table)
        return tables
    
    def _copy_slide_to_presentation(self, target_prs: Presentation, source_slide):
        """将源幻灯片复制到目标PPT的末尾"""
        # 使用源幻灯片的布局
        source_layout = source_slide.slide_layout
        
        # 在目标PPT中查找相同的布局
        target_layout = target_prs.slide_layouts[0]  # 默认使用第一个布局
        
        # 添加新幻灯片
        new_slide = target_prs.slides.add_slide(target_layout)
        
        # 复制所有形状
        for shape in source_slide.shapes:
            el = shape.element
            newel = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        return new_slide
    
    def _insert_slide_at_position(self, target_prs: Presentation, source_slide, position: int):
        """在指定位置插入幻灯片"""
        # 先添加到末尾
        new_slide = self._copy_slide_to_presentation(target_prs, source_slide)
        
        # 获取幻灯片的XML元素
        slides = list(target_prs.slides._sldIdLst)
        new_slide_element = slides[-1]
        
        # 移动到指定位置
        target_prs.slides._sldIdLst.remove(new_slide_element)
        target_prs.slides._sldIdLst.insert(position, new_slide_element)
        
        return position


if __name__ == "__main__":
    # 测试代码
    import json
    
    template_path = "副本人才引进ppt.pptx"
    data_path = "formatted_resumes.json"
    
    with open(data_path, 'r', encoding='utf-8') as f:
        resumes = json.load(f)
    
    renderer = OptimizedTemplateRendererV2(template_path)
    output_path = renderer.render_all(resumes)
    print(f"✅ 完成！输出: {output_path}")


# -*- coding: utf-8 -*-
"""
PPT 渲染
将结构化数据渲染到PowerPoint 演示文稿。
论文表格分页、内容布局、多份简历合并。

    - 封面页: 姓名、学历、专业
    - 个人信息页: 基本信息、教育背景、成果统计
    - 论文页: 表格展示，自动分页
    - 项目页: 获批项目情况
    - 获奖页: 获奖、成果、著作

"""

import os
import json
import math
from datetime import datetime
from typing import List, Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.table import Table


class ColorScheme:
    """
    配色方案
    
    采用学术风格配色，主色调为深蓝色，搭配金色强调。
    可根据需要修改这里的颜色值来自定义 PPT 风格。
    """
    PRIMARY = RgbColor(0x0F, 0x4C, 0x81)
    PRIMARY_LIGHT = RgbColor(0x1E, 0x88, 0xE5)
    PRIMARY_DARK = RgbColor(0x0A, 0x2F, 0x51)
    ACCENT_GOLD = RgbColor(0xD4, 0xAF, 0x37)
    ACCENT_TEAL = RgbColor(0x00, 0x96, 0x88)
    ACCENT_PURPLE = RgbColor(0x6A, 0x1B, 0x9A)
    ACCENT_CORAL = RgbColor(0xE6, 0x55, 0x50)
    DARK = RgbColor(0x1A, 0x1A, 0x2E)
    TEXT = RgbColor(0x4A, 0x4A, 0x5A)
    TEXT_LIGHT = RgbColor(0x6B, 0x7B, 0x8C)
    BG_LIGHT = RgbColor(0xF7, 0xF9, 0xFC)
    WHITE = RgbColor(0xFF, 0xFF, 0xFF)
    BORDER = RgbColor(0xE1, 0xE5, 0xEB)
    TABLE_HEADER = RgbColor(0x0F, 0x4C, 0x81)
    TABLE_ROW_ALT = RgbColor(0xF0, 0xF4, 0xF8)


class PaperItem:
    """
    单篇论文的数据结构
    
    用于表格展示，包含论文的所有展示字段。
    """
    def __init__(self, index: int, title: str, journal: str, category: str, 
                 year: str, is_first_author: str = "待确认"):
        self.index = index              # 序号
        self.title = title              # 论文标题
        self.journal = journal          # 期刊/会议名称
        self.category = category or "-" # 类别（SCI/CCF等）
        self.year = year                # 发表年份
        self.is_first_author = is_first_author  # 是否一作/通讯


class DataFormatter:
    """
    数据格式化工具
    
    将 API 返回的原始数据转换为 PPT 展示所需的格式。
    处理字段缺失、格式不统一等问题。
    """
    
    @staticmethod
    def format_education(edu_list: List[Dict]) -> str:
        """格式化教育背景"""
        if not edu_list:
            return "暂无教育背景信息"
        
        lines = []
        for edu in edu_list:
            time_range = edu.get("时间区间", "")
            school = edu.get("院校", "")
            major = edu.get("专业", "")
            degree = edu.get("学位", "")
            
            line = f"• {time_range}  {school}"
            if major:
                line += f"  {major}"
            if degree:
                line += f"  ({degree})"
            lines.append(line)
        
        return "\n".join(lines)
    
    @staticmethod
    def extract_paper_items(papers: List[Dict]) -> List[PaperItem]:
        """
        将论文数据展开为单篇论文列表
        
        Args:
            papers: 原始论文数据（按期刊分组）
            
        Returns:
            List[PaperItem]: 展开后的论文列表
        """
        items = []
        index = 1
        
        for paper_group in papers:
            journal = paper_group.get("期刊名称", "").replace("《", "").replace("》", "")
            category = paper_group.get("类别", "") or "-"
            years = paper_group.get("年份", "") or paper_group.get("年 份", "") or paper_group.get(" 年份", "")
            titles = paper_group.get("论文题目列表", [])
            
            # 将年份分割，尝试与论文标题对应
            year_list = [y.strip() for y in str(years).replace("、", ",").split(",") if y.strip()]
            
            for i, title in enumerate(titles):
                # 尝试匹配年份（如果年份数量与论文数量一致）
                year = year_list[i] if i < len(year_list) else (year_list[0] if year_list else "-")
                
                items.append(PaperItem(
                    index=index,
                    title=title,
                    journal=journal,
                    category=category,
                    year=year,
                    is_first_author="待确认"  # 数据中暂无此信息，后续可扩展
                ))
                index += 1
        
        return items
    
    @staticmethod
    def format_papers_summary(papers: List[Dict]) -> str:
        """生成论文统计摘要"""
        if not papers:
            return "暂无论文发表记录"
        
        total_papers = 0
        sci_count = 0
        sci1_count = 0
        ccf_a_count = 0
        
        for paper in papers:
            count_str = paper.get("篇数", "0") or paper.get(" 篇数", "0")
            count = int(''.join(filter(str.isdigit, str(count_str))) or 0)
            total_papers += count
            
            category = paper.get("类别", "") or ""
            if "SCI" in category:
                sci_count += count
                if "1区" in category:
                    sci1_count += count
            if "CCF A" in category.upper():
                ccf_a_count += count
        
        summary_parts = [f"共发表论文 {total_papers} 篇"]
        if sci1_count > 0:
            summary_parts.append(f"SCI 1区 {sci1_count} 篇")
        elif sci_count > 0:
            summary_parts.append(f"SCI {sci_count} 篇")
        if ccf_a_count > 0:
            summary_parts.append(f"CCF A类 {ccf_a_count} 篇")
        
        return " | ".join(summary_parts)
    
    @staticmethod
    def format_projects(projects: List[Dict]) -> str:
        """格式化项目列表"""
        if not projects:
            return "暂无项目信息"
        
        lines = []
        for proj in projects:
            category = proj.get("项目类别", "") or proj.get("项 目类别", "")
            count = proj.get("项数", "")
            years = proj.get("年份", "")
            budget = proj.get("备注", "") or proj.get("备 注", "")
            names = proj.get("项目名称列表", []) or proj.get("项目名称列 表", [])
            
            header = f"【{category}】 {count}"
            if years:
                header += f" ({years})"
            if budget:
                header += f" - 经费: {budget}"
            lines.append(header)
            
            if names:
                for name in names:
                    lines.append(f"   • {name}")
            
            lines.append("")
        
        return "\n".join(lines).strip()
    
    @staticmethod
    def format_awards(awards: List[Dict]) -> str:
        """格式化获奖情况"""
        if not awards:
            return "暂无获奖记录"
        
        lines = []
        for award in awards:
            name = award.get("奖项名称", "")
            year = award.get("年份", "")
            award_type = award.get("类型", "") or award.get("类 型", "")
            
            line = f"🏅 {name}"
            if year:
                line += f" ({year}年)"
            if award_type:
                line += f" - {award_type}"
            lines.append(line)
        
        return "\n".join(lines)
    
    @staticmethod
    def format_other_achievements(achievements: List[Dict]) -> str:
        """格式化其他成果"""
        if not achievements:
            return "暂无其他成果"
        
        lines = []
        for ach in achievements:
            category = ach.get("类别", "")
            count = ach.get("项数", "")
            years = ach.get("年份", "")
            note = ach.get("备注", "")
            names = ach.get("名称列表", [])
            
            line = f"• {category}: {count}"
            if years:
                line += f" ({years})"
            if note:
                line += f" - {note}"
            lines.append(line)
            
            if names:
                for name in names[:3]:
                    lines.append(f"   - {name}")
        
        return "\n".join(lines)
    
    @staticmethod
    def format_publications(publications: Any) -> str:
        """格式化著作情况"""
        if not publications:
            return "暂无著作出版"
        
        if isinstance(publications, list):
            lines = []
            for pub in publications:
                name = pub.get("著作名称", "")
                publisher = pub.get("出版社", "") or ""
                time = pub.get("出版时间", "") or ""
                
                line = f"📖 {name}"
                if time:
                    line += f" ({time})"
                if publisher:
                    line += f" - {publisher}"
                lines.append(line)
            return "\n".join(lines)
        
        return str(publications)
    
    @staticmethod
    def count_total_papers(papers: List[Dict]) -> int:
        """统计论文总数"""
        total = 0
        for paper in papers:
            count_str = paper.get("篇数", "0") or paper.get(" 篇数", "0")
            count = int(''.join(filter(str.isdigit, str(count_str))) or 0)
            total += count
        return total
    
    @staticmethod
    def count_total_projects(projects: List[Dict]) -> int:
        """统计项目总数"""
        total = 0
        for proj in projects:
            count_str = proj.get("项数", "0")
            count = int(''.join(filter(str.isdigit, str(count_str))) or 0)
            total += count
        return total


class PPTRenderer:
 
    #PPT 渲染引擎
   
    # 论文表格布局配置
    PAPERS_PER_PAGE = 8                 # 每页论文数量
    TABLE_HEADER_HEIGHT = Inches(0.4)   # 表头高度
    TABLE_ROW_HEIGHT = Inches(0.55)     # 数据行高度
    
    # 表格列定义: (列名, 列宽)
    TABLE_COLUMNS = [
        ("序号", Inches(0.5)),
        ("论文题目", Inches(5.2)),
        ("期刊/会议", Inches(2.8)),
        ("类别", Inches(1.0)),
        ("年份", Inches(0.7)),
        ("一作/通讯", Inches(0.8)),
    ]
    
    def __init__(self, papers_per_page: int = 8):
        """
        初始化渲染器
        
        Args:
            papers_per_page: 每页论文数量，根据论文标题长度可调整为 6-10
        """
        self.formatter = DataFormatter()
        self.slide_width = Inches(13.333)   # 16:9 宽屏
        self.slide_height = Inches(7.5)
        self.PAPERS_PER_PAGE = papers_per_page
    
    def _add_multiline_text(self, text_frame, content: str, font_size: int = 12, 
                            color: RgbColor = None, line_spacing: float = 1.15):
        """添加多行文本到文本框"""
        lines = content.split('\n')
        
        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line
            p.font.size = Pt(font_size)
            p.font.name = "Microsoft YaHei"
            if color:
                p.font.color.rgb = color
            else:
                p.font.color.rgb = ColorScheme.TEXT
            
            p.line_spacing = line_spacing
    
    def _add_background(self, slide, color=None):
        """设置幻灯片背景"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color or ColorScheme.BG_LIGHT
    
    def _add_top_banner(self, slide, color, title, icon="", subtitle=""):
        """添加顶部横幅"""
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.slide_width, Inches(1.1)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = color
        banner.line.fill.background()
        
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.1),
            self.slide_width, Inches(0.05)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = ColorScheme.ACCENT_GOLD
        stripe.line.fill.background()
        
        title_text = f"{icon} {title}" if icon else title
        if subtitle:
            title_text += f"  {subtitle}"
        
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.25),
            Inches(12), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ColorScheme.WHITE
        p.font.name = "Microsoft YaHei"
    
    def _add_content_card(self, slide, left, top, width, height):
        """添加内容卡片"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top,
            width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = ColorScheme.WHITE
        card.line.color.rgb = ColorScheme.BORDER
        card.line.width = Pt(1)
        
        try:
            card.adjustments[0] = 0.05
        except:
            pass
        
        return card
    
    def _add_section_title(self, slide, title, left, top, width, icon=""):
        """添加区块标题"""
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top,
            width, Inches(0.45)
        )
        container.fill.solid()
        container.fill.fore_color.rgb = ColorScheme.PRIMARY
        container.line.fill.background()
        
        try:
            container.adjustments[0] = 0.15
        except:
            pass
        
        tf = container.text_frame
        tf.paragraphs[0].text = f"{icon} {title}" if icon else title
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = ColorScheme.WHITE
        tf.paragraphs[0].font.name = "Microsoft YaHei"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return container
    
    def _add_text_box(self, slide, content, left, top, width, height, 
                      font_size=12, color=None, bold=False):
        """添加文本框"""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        
        self._add_multiline_text(tf, content, font_size, color)
        
        if bold:
            for p in tf.paragraphs:
                p.font.bold = True
        
        return text_box
    
    def _set_cell_style(self, cell, text: str, font_size: int = 9, bold: bool = False,
                        color: RgbColor = None, bg_color: RgbColor = None,
                        alignment: PP_ALIGN = PP_ALIGN.LEFT):
        """设置表格单元格样式"""
        cell.text = text
        
        # 设置文本框属性
        tf = cell.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        # 设置段落属性
        for p in tf.paragraphs:
            p.font.size = Pt(font_size)
            p.font.name = "Microsoft YaHei"
            p.font.bold = bold
            p.alignment = alignment
            if color:
                p.font.color.rgb = color
        
        # 设置背景色
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
    
    def _create_papers_table(self, slide, papers: List[PaperItem], 
                             left: float, top: float, width: float) -> Table:
        """
        创建论文表格（使用固定行高）
        
        Args:
            slide: 幻灯片对象
            papers: 论文列表
            left, top, width: 表格位置和宽度
            
        Returns:
            Table: 表格对象
        """
        rows = len(papers) + 1  # 数据行 + 表头
        cols = len(self.TABLE_COLUMNS)
        
        # 使用固定行高计算表格总高度
        table_height = self.TABLE_HEADER_HEIGHT + self.TABLE_ROW_HEIGHT * len(papers)
        
        # 创建表格
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, table_height)
        table = table_shape.table
        
        # 设置固定行高
        table.rows[0].height = self.TABLE_HEADER_HEIGHT
        for i in range(1, rows):
            table.rows[i].height = self.TABLE_ROW_HEIGHT
        
        # 设置列宽
        for i, (_, col_width) in enumerate(self.TABLE_COLUMNS):
            table.columns[i].width = col_width
        
        # 设置表头
        headers = [col[0] for col in self.TABLE_COLUMNS]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            self._set_cell_style(
                cell, header,
                font_size=10, bold=True,
                color=ColorScheme.WHITE,
                bg_color=ColorScheme.TABLE_HEADER,
                alignment=PP_ALIGN.CENTER
            )
        
        # 填充数据行
        for row_idx, paper in enumerate(papers):
            row = row_idx + 1  # 跳过表头
            is_alt_row = row_idx % 2 == 1
            bg_color = ColorScheme.TABLE_ROW_ALT if is_alt_row else ColorScheme.WHITE
            
            # 序号
            self._set_cell_style(
                table.cell(row, 0), str(paper.index),
                font_size=9, bg_color=bg_color, alignment=PP_ALIGN.CENTER
            )
            
            # 论文题目（较长，可能需要截断）
            title = paper.title
            if len(title) > 80:
                title = title[:77] + "..."
            self._set_cell_style(
                table.cell(row, 1), title,
                font_size=8, bg_color=bg_color, alignment=PP_ALIGN.LEFT
            )
            
            # 期刊/会议名称
            journal = paper.journal
            if len(journal) > 40:
                journal = journal[:37] + "..."
            self._set_cell_style(
                table.cell(row, 2), journal,
                font_size=8, bg_color=bg_color, alignment=PP_ALIGN.LEFT
            )
            
            # 类别
            self._set_cell_style(
                table.cell(row, 3), paper.category,
                font_size=8, bg_color=bg_color, alignment=PP_ALIGN.CENTER
            )
            
            # 年份
            self._set_cell_style(
                table.cell(row, 4), paper.year,
                font_size=8, bg_color=bg_color, alignment=PP_ALIGN.CENTER
            )
            
            # 是否一作/第一通讯
            self._set_cell_style(
                table.cell(row, 5), paper.is_first_author,
                font_size=8, bg_color=bg_color, alignment=PP_ALIGN.CENTER
            )
        
        return table
    
    def render_cover_slide(self, prs: Presentation, resume: Dict) -> None:
        """渲染封面幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide, RgbColor(0xFD, 0xFB, 0xF7))
        
        # 顶部装饰带
        top_band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.slide_width, Inches(0.12)
        )
        top_band.fill.solid()
        top_band.fill.fore_color.rgb = ColorScheme.PRIMARY
        top_band.line.fill.background()
        
        # 左侧装饰块
        left_block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.4), self.slide_height
        )
        left_block.fill.solid()
        left_block.fill.fore_color.rgb = ColorScheme.PRIMARY_DARK
        left_block.line.fill.background()
        
        # 装饰金线
        gold_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4), Inches(0),
            Inches(0.03), self.slide_height
        )
        gold_line.fill.solid()
        gold_line.fill.fore_color.rgb = ColorScheme.ACCENT_GOLD
        gold_line.line.fill.background()
        
        # 右上角装饰圆形
        circle1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            self.slide_width - Inches(2.5), Inches(-0.8),
            Inches(3.5), Inches(3.5)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = ColorScheme.PRIMARY
        circle1.fill.fore_color.brightness = 0.85
        circle1.line.fill.background()
        
        # 右下角装饰
        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            self.slide_width - Inches(1.8), self.slide_height - Inches(1.2),
            Inches(2.2), Inches(2.2)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = ColorScheme.ACCENT_TEAL
        circle2.fill.fore_color.brightness = 0.7
        circle2.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.3),
            Inches(10), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "人才引进简历"
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = ColorScheme.DARK
        p.font.name = "Microsoft YaHei"
        
        # 英文副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(3.4),
            Inches(10), Inches(0.6)
        )
        tf2 = subtitle_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "TALENT RECRUITMENT RESUME"
        p2.font.size = Pt(18)
        p2.font.color.rgb = ColorScheme.TEXT_LIGHT
        p2.font.name = "Arial"
        p2.font.bold = True
        
        # 装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.2), Inches(4.1),
            Inches(4), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ColorScheme.ACCENT_GOLD
        line.line.fill.background()
        
        # 姓名
        name = resume.get("姓名", "未知")
        name_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(4.5),
            Inches(8), Inches(0.9)
        )
        tf3 = name_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = name
        p3.font.size = Pt(42)
        p3.font.bold = True
        p3.font.color.rgb = ColorScheme.PRIMARY
        p3.font.name = "Microsoft YaHei"
        
        # 学历信息
        degree = resume.get("最高学历", "")
        edu_list = resume.get("就读院校", [])
        field = ""
        if edu_list:
            first_edu = edu_list[0]
            field = first_edu.get("专业", "")
        
        info_text = degree
        if field:
            info_text += f" · {field}"
        
        degree_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(5.5),
            Inches(8), Inches(0.5)
        )
        tf4 = degree_box.text_frame
        p4 = tf4.paragraphs[0]
        p4.text = info_text
        p4.font.size = Pt(18)
        p4.font.color.rgb = ColorScheme.TEXT
        p4.font.name = "Microsoft YaHei"
    
    def render_profile_slide(self, prs: Presentation, resume: Dict) -> None:
        """渲染个人信息幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide)
        self._add_top_banner(slide, ColorScheme.PRIMARY, "个人基本信息", "👤")
        
        # 左侧：基本信息卡片
        self._add_content_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.4))
        
        # 姓名
        name = resume.get("姓名", "未知")
        name_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6), Inches(5), Inches(0.7)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ColorScheme.DARK
        p.font.name = "Microsoft YaHei"
        
        # 基本信息
        info_items = [
            ("性别", resume.get("性别", "未知")),
            ("年龄", f"{resume.get('年龄', '未知')}岁"),
            ("出生年份", f"{resume.get('出生日期', '未知')}年"),
            ("最高学历", resume.get("最高学历", "未知")),
        ]
        
        for i, (label, value) in enumerate(info_items):
            col = i % 2
            row = i // 2
            left = Inches(0.8 + col * 2.8)
            top = Inches(2.3 + row * 0.55)
            
            label_box = slide.shapes.add_textbox(left, top, Inches(1), Inches(0.35))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{label}："
            p.font.size = Pt(13)
            p.font.color.rgb = ColorScheme.TEXT_LIGHT
            p.font.name = "Microsoft YaHei"
            
            value_box = slide.shapes.add_textbox(left + Inches(1), top, Inches(1.6), Inches(0.35))
            tf2 = value_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = str(value) if value and value != "未知" else "未知"
            p2.font.size = Pt(13)
            p2.font.bold = True
            p2.font.color.rgb = ColorScheme.DARK
            p2.font.name = "Microsoft YaHei"
        
        # 右侧统计卡片
        papers_count = self.formatter.count_total_papers(resume.get("发表论文情况", []))
        projects_count = self.formatter.count_total_projects(resume.get("获批项目情况", []))
        awards_count = len(resume.get("获奖情况", []))
        
        stats = [
            ("📄", "论文", str(papers_count), ColorScheme.PRIMARY_LIGHT),
            ("💼", "项目", str(projects_count), ColorScheme.ACCENT_TEAL),
            ("🏆", "获奖", str(awards_count), ColorScheme.ACCENT_GOLD),
        ]
        
        for i, (icon, label, value, color) in enumerate(stats):
            left = Inches(6.6 + i * 2.1)
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, Inches(1.4),
                Inches(1.9), Inches(1.1)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = color
            card.line.fill.background()
            
            stat_box = slide.shapes.add_textbox(left, Inches(1.5), Inches(1.9), Inches(0.9))
            tf = stat_box.text_frame
            
            p1 = tf.paragraphs[0]
            p1.text = value
            p1.font.size = Pt(24)
            p1.font.bold = True
            p1.font.color.rgb = ColorScheme.WHITE
            p1.font.name = "Arial"
            p1.alignment = PP_ALIGN.CENTER
            
            p2 = tf.add_paragraph()
            p2.text = label
            p2.font.size = Pt(11)
            p2.font.color.rgb = ColorScheme.WHITE
            p2.font.name = "Microsoft YaHei"
            p2.alignment = PP_ALIGN.CENTER
        
        # 教育背景
        self._add_section_title(slide, "教育背景", Inches(0.5), Inches(4.0), Inches(12.3), "🎓")
        self._add_content_card(slide, Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.6))
        
        edu_text = self.formatter.format_education(resume.get("就读院校", []))
        self._add_text_box(slide, edu_text, Inches(0.8), Inches(4.75), Inches(11.8), Inches(2.2), font_size=13)
    
    def render_papers_slides(self, prs: Presentation, resume: Dict) -> int:
        
        #渲染论文成果幻灯片（分页）
      
        papers = resume.get("发表论文情况", [])
        
        if not papers:
            # 如果没有论文，创建一个空页面
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_background(slide)
            self._add_top_banner(slide, ColorScheme.ACCENT_PURPLE, "发表论文情况", "📚")
            
            self._add_content_card(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.75))
            self._add_text_box(slide, "暂无论文发表记录", Inches(0.8), Inches(3.5), 
                              Inches(11.8), Inches(1), font_size=16)
            return 1
        
        # 展开论文为单篇列表
        paper_items = self.formatter.extract_paper_items(papers)
        total_papers = len(paper_items)
        
        # 计算需要多少页
        total_pages = math.ceil(total_papers / self.PAPERS_PER_PAGE)
        
        # 生成论文统计摘要
        summary = self.formatter.format_papers_summary(papers)
        
        # 为每页创建幻灯片
        for page_num in range(total_pages):
            start_idx = page_num * self.PAPERS_PER_PAGE
            end_idx = min(start_idx + self.PAPERS_PER_PAGE, total_papers)
            page_papers = paper_items[start_idx:end_idx]
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_background(slide)
            
            # 页面标题（包含页码）
            page_info = f"（第 {page_num + 1}/{total_pages} 页，共 {total_papers} 篇）"
            self._add_top_banner(slide, ColorScheme.ACCENT_PURPLE, "发表论文情况", "📚", page_info)
            
            # 论文摘要（仅第一页显示完整摘要）
            if page_num == 0:
                self._add_content_card(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.55))
                self._add_text_box(slide, summary, Inches(0.8), Inches(1.4), Inches(11.8), Inches(0.4), 
                                  font_size=13, color=ColorScheme.DARK, bold=True)
                table_top = Inches(2.0)
            else:
                table_top = Inches(1.3)
            
            # 创建论文表格（使用固定行高，自动计算高度）
            self._create_papers_table(
                slide, page_papers,
                left=Inches(0.5),
                top=table_top,
                width=Inches(12.3)
            )
        
        return total_pages
    
    def render_projects_slide(self, prs: Presentation, resume: Dict) -> None:
        """渲染项目情况幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide)
        self._add_top_banner(slide, ColorScheme.ACCENT_TEAL, "获批项目情况", "💼")
        
        self._add_content_card(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.75))
        
        projects = resume.get("获批项目情况", [])
        projects_text = self.formatter.format_projects(projects)
        self._add_text_box(slide, projects_text, Inches(0.8), Inches(1.6), Inches(11.8), Inches(5.35), font_size=12)
    
    def render_awards_slide(self, prs: Presentation, resume: Dict) -> None:
        """渲染获奖与成果幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide)
        self._add_top_banner(slide, ColorScheme.ACCENT_GOLD, "获奖与成果", "🏆")
        
        # 获奖情况
        self._add_section_title(slide, "获奖情况", Inches(0.5), Inches(1.4), Inches(12.3), "🎖️")
        self._add_content_card(slide, Inches(0.5), Inches(1.95), Inches(12.3), Inches(2.0))
        
        awards = resume.get("获奖情况", [])
        awards_text = self.formatter.format_awards(awards)
        self._add_text_box(slide, awards_text, Inches(0.8), Inches(2.15), Inches(11.8), Inches(1.6), font_size=13)
        
        # 其他成果
        self._add_section_title(slide, "其他成果", Inches(0.5), Inches(4.15), Inches(6.0), "🔬")
        self._add_content_card(slide, Inches(0.5), Inches(4.7), Inches(6.0), Inches(2.45))
        
        other = resume.get("其他成果", [])
        other_text = self.formatter.format_other_achievements(other)
        self._add_text_box(slide, other_text, Inches(0.8), Inches(4.9), Inches(5.4), Inches(2.05), font_size=11)
        
        # 著作情况
        self._add_section_title(slide, "著作情况", Inches(6.8), Inches(4.15), Inches(6.0), "📖")
        self._add_content_card(slide, Inches(6.8), Inches(4.7), Inches(6.0), Inches(2.45))
        
        publications = resume.get("著作情况", [])
        pub_text = self.formatter.format_publications(publications)
        self._add_text_box(slide, pub_text, Inches(7.1), Inches(4.9), Inches(5.4), Inches(2.05), font_size=11)
    
    def render_resume(self, resume: Dict) -> Presentation:
        """
        渲染单份简历为PPT
        
        Args:
            resume: 格式化后的简历数据
            
        Returns:
            Presentation 对象
        """
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        self.render_cover_slide(prs, resume)
        self.render_profile_slide(prs, resume)
        self.render_papers_slides(prs, resume)  # 使用新的分页方法
        self.render_projects_slide(prs, resume)
        self.render_awards_slide(prs, resume)
        
        return prs
    
    def render_all(self, resumes: List[Dict], output_dir: str = None) -> str:
        #渲染多份简历到一个PPT
       
        if not resumes:
            raise ValueError("简历列表不能为空")
        
        if output_dir is None:
            output_dir = os.path.dirname(__file__)
        
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        for i, resume in enumerate(resumes):
            name = resume.get('姓名', '未知')
            papers_count = self.formatter.count_total_papers(resume.get("发表论文情况", []))
            print(f"📝 正在渲染第 {i+1} 份简历: {name} (论文 {papers_count} 篇)")
            
            self.render_cover_slide(prs, resume)
            self.render_profile_slide(prs, resume)
            
            # 论文分页渲染
            paper_pages = self.render_papers_slides(prs, resume)
            print(f"   📄 论文页面: {paper_pages} 页")
            
            self.render_projects_slide(prs, resume)
            self.render_awards_slide(prs, resume)
            
            # 添加分隔页（如果不是最后一份）
            if i < len(resumes) - 1:
                self._add_separator_slide(prs, i + 2, len(resumes))
        
        # 生成输出文件名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(output_dir, f"人才引进简历_{timestamp}.pptx")
        
        prs.save(output_path)
        print(f"✅ PPT已保存: {output_path}")
        
        return output_path
    
    def _add_separator_slide(self, prs: Presentation, next_num: int, total: int):
        """添加分隔页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide, ColorScheme.PRIMARY_DARK)
        
        # 中心文字
        text_box = slide.shapes.add_textbox(
            Inches(0), Inches(3),
            self.slide_width, Inches(1.5)
        )
        tf = text_box.text_frame
        
        p = tf.paragraphs[0]
        p.text = f"第 {next_num} / {total} 份简历"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ColorScheme.WHITE
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER


def main():
    """主函数 - 测试渲染器"""
    # 尝试从文件加载数据
    data_file = os.path.join(os.path.dirname(__file__), "formatted_resumes.json")
    
    if os.path.exists(data_file):
        print(f"📂 从文件加载数据: {data_file}")
        with open(data_file, "r", encoding="utf-8") as f:
            resumes = json.load(f)
    else:
        # 从 response.txt 解析
        response_file = os.path.join(os.path.dirname(__file__), "response.txt")
        if os.path.exists(response_file):
            print(f"📂 从响应文件解析: {response_file}")
            from run import parse_from_response_file
            resumes = parse_from_response_file(response_file)
        else:
            print("❌ 未找到数据文件，请先运行 run.py 获取数据")
            return
    
    if resumes:
        # 创建渲染器，可以调整每页论文数量
        renderer = PPTRenderer(papers_per_page=8)
        output_path = renderer.render_all(resumes)
        print(f"\n🎉 PPT生成完成: {output_path}")
    else:
        print("❌ 无法加载简历数据")


if __name__ == "__main__":
    main()
